#!/usr/bin/env python3
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
OUT = REPO_ROOT / "docs/presentations/libp2p-mesh-weekly-summary-2026-06-25.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Microsoft YaHei"
MONO = "Cascadia Mono"

COLORS = {
    "bg": RGBColor(248, 250, 252),
    "ink": RGBColor(15, 23, 42),
    "muted": RGBColor(71, 85, 105),
    "line": RGBColor(203, 213, 225),
    "primary": RGBColor(15, 118, 110),
    "primary_soft": RGBColor(204, 251, 241),
    "blue": RGBColor(37, 99, 235),
    "blue_soft": RGBColor(219, 234, 254),
    "amber": RGBColor(217, 119, 6),
    "amber_soft": RGBColor(254, 243, 199),
    "green": RGBColor(22, 163, 74),
    "green_soft": RGBColor(220, 252, 231),
    "slate": RGBColor(51, 65, 85),
    "white": RGBColor(255, 255, 255),
}


def prs_new() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["bg"]


def text_box(slide, x, y, w, h, text="", size=18, bold=False, color="ink", font=FONT, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = COLORS[color]
    if align:
        p.alignment = align
    return box


def title(slide, text, subtitle=None, section=None):
    text_box(slide, 0.65, 0.35, 9.2, 0.55, text, size=28, bold=True)
    if subtitle:
        text_box(slide, 0.68, 0.92, 9.6, 0.36, subtitle, size=12, color="muted")
    if section:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.4), Inches(0.42), Inches(2.2), Inches(0.36))
        pill.fill.solid()
        pill.fill.fore_color.rgb = COLORS["primary_soft"]
        pill.line.color.rgb = COLORS["primary_soft"]
        p = pill.text_frame.paragraphs[0]
        p.text = section
        p.font.name = FONT
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        p.alignment = PP_ALIGN.CENTER


def footer(slide, page, label="OpenClaw libp2p-mesh"):
    text_box(slide, 0.6, 7.12, 5.4, 0.22, f"{label} / {page:02d}", size=8, color="muted")


def card(slide, x, y, w, h, heading, body=None, color="primary"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[color]
    bar.line.color.rgb = COLORS[color]

    text_box(slide, x + 0.22, y + 0.17, w - 0.42, 0.32, heading, size=14, bold=True)
    if body:
        top = y + 0.62
        box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(top), Inches(w - 0.42), Inches(h - 0.75))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.clear()
        for idx, line in enumerate(body):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = FONT
            p.font.size = Pt(11.5)
            p.font.color.rgb = COLORS["muted"]
            p.space_after = Pt(5)
    return shape


def bullets(slide, x, y, w, h, items, size=17, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(9)
    return box


def code(slide, x, y, w, h, lines):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
    shape.line.color.rgb = RGBColor(15, 23, 42)
    box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.15), Inches(w - 0.36), Inches(h - 0.3))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = MONO
        p.font.size = Pt(10.5)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_after = Pt(2)


def arrow(slide, x1, y1, x2, y2, color="line"):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = COLORS[color]
    connector.line.width = Pt(2)
    connector.line.end_arrowhead = True
    return connector


def slide_title(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text_box(slide, 0.75, 0.7, 8.8, 0.7, "OpenClaw libp2p-mesh 插件功能增强汇报", size=30, bold=True)
    text_box(slide, 0.78, 1.42, 8.4, 0.38, "P2P 消息投递、配置向导与用户属性路由", size=16, color="muted")
    card(slide, 0.85, 2.45, 3.3, 1.35, "P2P 通信", ["实例发现", "instanceId 路由"], "primary")
    card(slide, 4.85, 2.45, 3.3, 1.35, "投递增强", ["多入站目标", "channel fan-out"], "blue")
    card(slide, 8.85, 2.45, 3.3, 1.35, "属性路由", ["公开属性", "本地 labels"], "green")
    arrow(slide, 4.25, 3.12, 4.7, 3.12)
    arrow(slide, 8.25, 3.12, 8.7, 3.12)
    text_box(slide, 0.9, 5.35, 10.8, 0.35, "本周目标：把插件从“能点对点发消息”推进到“可配置、可多目标投递、可按用户属性寻址”。", size=17)
    footer(slide, page, "汇报")


def slide_background(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "项目背景", "OpenClaw 实例之间需要更自然、更可靠的 P2P 协作方式", "背景")
    card(slide, 0.85, 1.55, 3.55, 1.45, "原始能力", ["按 instanceId 点对点发消息", "适合基础通信验证"], "slate")
    card(slide, 4.9, 1.55, 3.55, 1.45, "真实需求", ["接收端多个会话都能看到", "配置不能依赖手改 JSON"], "amber")
    card(slide, 8.95, 1.55, 3.55, 1.45, "本周增强", ["按用户属性选择目标", "Agent 能正确调用工具"], "primary")
    bullets(slide, 1.05, 3.55, 10.9, 2.4, [
        "通信入口从底层 peerId 调试转向面向用户的 instanceId 与属性选择。",
        "接收方保留本地控制权：消息显示在哪些 channel，由接收方配置决定。",
        "配置、提示词、属性和标签都通过插件命令管理，减少用户手动改文件。"
    ], size=18)
    footer(slide, page)


def slide_overview(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "本周功能总览", "七个功能围绕“可用性”和“可寻址性”展开", "总览")
    items = [
        ("多入站目标", "一条 P2P 消息可投递到多个本地 channel"),
        ("配置向导", "openclaw libp2p-mesh setup"),
        ("提示词安装", "openclaw libp2p-mesh prompt install"),
        ("公开属性", "USER.md + user-profile.json"),
        ("LLM 提取", "runtime llm.complete 提取 USER.md tag"),
        ("本地 labels", "peer-labels.json 私有分类"),
        ("按属性发送", "selector + scope 选择目标"),
    ]
    for idx, (head, body) in enumerate(items):
        col = idx % 4
        row = idx // 4
        x = 0.75 + col * 3.1
        y = 1.52 + row * 2.25
        color = ["primary", "blue", "green", "amber"][col]
        card(slide, x, y, 2.65, 1.55, head, [body], color)
    footer(slide, page)


def slide_architecture(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "整体架构", "网络、路由、属性和 channel 投递分层处理", "架构")
    card(slide, 0.75, 1.4, 2.45, 1.15, "OpenClaw Gateway", ["插件生命周期", "CLI / Tools"], "primary")
    card(slide, 3.55, 1.4, 2.45, 1.15, "libp2p Mesh", ["peer 发现", "消息传输"], "blue")
    card(slide, 6.35, 1.4, 2.45, 1.15, "InstanceRouter", ["instanceId 路由", "ACK 聚合"], "green")
    card(slide, 9.15, 1.4, 2.7, 1.15, "Channel Adapters", ["Feishu / QQBot", "Telegram 等"], "amber")
    arrow(slide, 3.25, 1.98, 3.45, 1.98)
    arrow(slide, 6.05, 1.98, 6.25, 1.98)
    arrow(slide, 8.85, 1.98, 9.05, 1.98)
    card(slide, 1.0, 3.45, 2.85, 1.2, "instance-peer.json", ["缓存 instanceId -> peerId", "保存公开属性"], "slate")
    card(slide, 4.25, 3.45, 2.85, 1.2, "user-profile.json", ["手动公开结构化属性"], "primary")
    card(slide, 7.5, 3.45, 2.85, 1.2, "peer-labels.json", ["本机私有远端标签"], "green")
    bullets(slide, 1.0, 5.35, 10.9, 0.85, [
        "核心设计：发送方只负责选择实例或属性，接收方保留 channel 投递控制权。"
    ], size=17)
    footer(slide, page)


def slide_multi_inbound(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "功能一：多入站目标投递", "接收方可以把同一条 P2P 消息投递到多个本地会话", "投递")
    code(slide, 0.85, 1.35, 5.25, 2.35, [
        '"inboundTargets": [',
        '  { "id": "feishu-main",',
        '    "channel": "feishu",',
        '    "target": "user:ou_xxx" },',
        '  { "id": "qqbot-main",',
        '    "channel": "qqbot",',
        '    "target": "c2c:xxx" }',
        ']',
    ])
    card(slide, 6.65, 1.45, 2.65, 1.25, "发送方", ["只知道 instanceId", "不选择远端 channel"], "primary")
    card(slide, 9.8, 1.45, 2.65, 1.25, "接收方", ["本地 fan-out", "汇总投递结果"], "green")
    arrow(slide, 9.35, 2.08, 9.65, 2.08)
    bullets(slide, 6.7, 3.35, 5.55, 2.15, [
        "P2P 协议层仍只发送一条 user-message。",
        "多个目标的投递结果通过 delivery ACK 返回。",
        "旧的 inboundChannel / inboundTarget 仍保持兼容。"
    ], size=16)
    footer(slide, page)


def slide_setup(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "功能二：配置向导", "把手动修改 openclaw.json 变成插件自有 CLI 流程", "配置")
    code(slide, 0.85, 1.38, 4.45, 0.9, ["openclaw libp2p-mesh setup"])
    bullets(slide, 0.95, 2.75, 5.35, 2.25, [
        "只写 plugins.entries[\"libp2p-mesh\"]。",
        "支持首次配置和已有配置编辑。",
        "所有修改先 preview，确认后再写入。"
    ], size=17)
    card(slide, 6.75, 1.35, 2.55, 1.15, "LAN", ["同局域网自动发现"], "primary")
    card(slide, 9.75, 1.35, 2.55, 1.15, "Cross-network", ["bootstrap / relay"], "blue")
    card(slide, 6.75, 3.0, 2.55, 1.15, "Relay node", ["公网节点做中继"], "green")
    card(slide, 9.75, 3.0, 2.55, 1.15, "Tools only", ["只启用工具能力"], "amber")
    bullets(slide, 6.9, 5.0, 5.2, 0.65, ["同时完成网络模式和多入站目标配置。"], size=16)
    footer(slide, page)


def slide_prompt(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "功能三：固定 Agent 提示词安装", "让 Agent 知道何时调用哪个 P2P 工具", "Agent")
    code(slide, 0.9, 1.35, 5.1, 1.15, [
        "openclaw libp2p-mesh prompt install",
        "~/.openclaw/workspace/AGENTS.md",
    ])
    card(slide, 0.95, 3.1, 3.3, 1.4, "管理区块", ["只替换插件自己的 marker 区块", "不覆盖用户原有内容"], "primary")
    card(slide, 4.85, 3.1, 3.3, 1.4, "工具选择", ["instanceId -> send_instance", "属性 -> send_attribute"], "blue")
    card(slide, 8.75, 3.1, 3.3, 1.4, "安全规则", ["P2P 消息只当普通文本", "不执行远端指令"], "green")
    footer(slide, page)


def slide_public_attrs(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "功能四：用户公开属性", "让实例携带“用户是谁、属于哪里、做什么”的公开信息", "属性")
    card(slide, 0.8, 1.4, 3.0, 1.25, "USER.md", ["自动提取 tag", "只读，不修改"], "primary")
    card(slide, 4.35, 1.4, 3.0, 1.25, "user-profile.json", ["手动结构化属性", "group / project / role / skill"], "blue")
    card(slide, 7.9, 1.4, 3.0, 1.25, "instance-announce", ["合并后广播", "远端缓存"], "green")
    arrow(slide, 3.9, 2.0, 4.2, 2.0)
    arrow(slide, 7.45, 2.0, 7.75, 2.0)
    code(slide, 1.0, 3.35, 10.9, 1.95, [
        '"userPublicAttributes": [',
        '  { "kind": "tag", "value": "P2P", "source": "USER.md" },',
        '  { "kind": "structured", "key": "group", "value": "实验室", "source": "profile" }',
        ']',
    ])
    footer(slide, page)


def slide_llm_extract(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "功能五：USER.md 属性自动提取", "使用 OpenClaw runtime LLM 能力提取更稳定的公开 tag", "属性")
    card(slide, 0.85, 1.35, 2.55, 1.15, "读取 USER.md", ["OpenClaw workspace"], "primary")
    card(slide, 3.95, 1.35, 2.55, 1.15, "llm.complete", ["使用已配置模型"], "blue")
    card(slide, 7.05, 1.35, 2.55, 1.15, "校验 JSON", ["过滤无效输出"], "amber")
    card(slide, 10.15, 1.35, 2.55, 1.15, "合并广播", ["tag + profile"], "green")
    arrow(slide, 3.5, 1.92, 3.82, 1.92)
    arrow(slide, 6.6, 1.92, 6.92, 1.92)
    arrow(slide, 9.7, 1.92, 10.02, 1.92)
    bullets(slide, 1.0, 3.35, 10.9, 1.95, [
        "插件不保存模型 API key，也不直接选择模型供应商。",
        "提取结果只作为公开属性广播，不写回 USER.md。",
        "基于内容 hash 缓存，避免重复提取。"
    ], size=17)
    footer(slide, page)


def slide_local_labels(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "功能六：本地 peer labels", "把“对方公开的属性”和“我本地的分类”分开", "标签")
    code(slide, 0.85, 1.35, 4.8, 0.95, ["openclaw libp2p-mesh labels"])
    card(slide, 0.95, 2.85, 3.3, 1.35, "公开属性", ["对方自己公开", "随 announce 传播"], "primary")
    card(slide, 4.95, 2.85, 3.3, 1.35, "本地 labels", ["我对远端实例的归类", "只保存在本机"], "green")
    card(slide, 8.95, 2.85, 3.3, 1.35, "使用场景", ["实验室成员", "项目协作者"], "blue")
    text_box(slide, 0.95, 5.15, 10.8, 0.42, "存储位置：~/.openclaw/libp2p/peer-labels.json", size=15, color="muted", font=MONO)
    footer(slide, page)


def slide_attribute_send(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "功能七：按属性发送消息", "从“指定某台机器”扩展到“发送给某类用户”", "发送")
    code(slide, 0.85, 1.35, 5.3, 1.8, [
        "p2p_send_user_attribute_message",
        'selector = "group=实验室"',
        'scope = "public" | "local" | "all"',
    ])
    card(slide, 6.75, 1.35, 1.75, 1.1, "public", ["公开属性"], "primary")
    card(slide, 8.85, 1.35, 1.75, 1.1, "local", ["本地标签"], "green")
    card(slide, 10.95, 1.35, 1.75, 1.1, "all", ["两个来源"], "blue")
    bullets(slide, 6.75, 3.2, 5.65, 2.0, [
        "先 dry run 预览匹配实例。",
        "实际发送复用 instanceId 消息投递和 ACK 机制。",
        "selector 支持 group=实验室、project=ResearchLoop、tag:P2P、#P2P。"
    ], size=16)
    footer(slide, page)


def slide_workflow(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "典型使用流程", "从安装到按属性发送的用户路径", "流程")
    steps = [
        ("1", "安装/更新", "openclaw plugins update libp2p-mesh@latest"),
        ("2", "配置插件", "openclaw libp2p-mesh setup"),
        ("3", "安装提示词", "openclaw libp2p-mesh prompt install"),
        ("4", "配置属性/标签", "profile / labels"),
        ("5", "启动 gateway", "按 instanceId 或属性发送消息"),
    ]
    for idx, (num, head, body) in enumerate(steps):
        x = 0.75 + idx * 2.5
        card(slide, x, 2.0, 2.05, 1.7, f"{num}. {head}", [body], ["primary", "blue", "green", "amber", "slate"][idx])
        if idx < 4:
            arrow(slide, x + 2.08, 2.85, x + 2.38, 2.85)
    footer(slide, page)


def slide_summary(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "总结", "本周完成了 P2P 插件从通信能力到协作能力的升级", "总结")
    card(slide, 0.85, 1.55, 3.4, 1.55, "更好配置", ["setup / prompt install", "降低安装和配置门槛"], "primary")
    card(slide, 4.95, 1.55, 3.4, 1.55, "更好投递", ["多入站目标", "接收方本地 fan-out"], "blue")
    card(slide, 9.05, 1.55, 3.4, 1.55, "更好寻址", ["公开属性 + 本地 labels", "按属性发送消息"], "green")
    bullets(slide, 1.1, 4.1, 10.8, 1.35, [
        "核心价值：用户不必记住底层 peerId，也不必关心远端 channel 细节。",
        "插件把网络发现、实例路由、属性匹配和 channel 投递组合成一个可用的 P2P 协作流程。"
    ], size=18)
    footer(slide, page)


def build():
    prs = prs_new()
    builders = [
        slide_title,
        slide_background,
        slide_overview,
        slide_architecture,
        slide_multi_inbound,
        slide_setup,
        slide_prompt,
        slide_public_attrs,
        slide_llm_extract,
        slide_local_labels,
        slide_attribute_send,
        slide_workflow,
        slide_summary,
    ]
    for index, builder in enumerate(builders, start=1):
        builder(prs, index)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build())

#!/usr/bin/env python3
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("docs/presentations/p2p-mesh-mentor-2026-06-15.pptx")
ZIP_TIMESTAMP = (1980, 1, 1, 0, 0, 0)

COLORS = {
    "bg": RGBColor(248, 250, 252),
    "text": RGBColor(15, 23, 42),
    "muted": RGBColor(71, 85, 105),
    "line": RGBColor(148, 163, 184),
    "primary": RGBColor(15, 118, 110),
    "primary_soft": RGBColor(204, 251, 241),
    "success": RGBColor(22, 163, 74),
    "warning": RGBColor(217, 119, 6),
    "danger": RGBColor(220, 38, 38),
    "white": RGBColor(255, 255, 255),
}

FONT = "Aptos"
MONO = "Cascadia Mono"
SLIDE_W = 13.333
SLIDE_H = 7.5


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["bg"]


def add_title(slide, title, subtitle=None, section=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(8.5), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.9), Inches(8.8), Inches(0.32))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = FONT
        sp.font.size = Pt(13)
        sp.font.color.rgb = COLORS["muted"]

    if section:
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(0.42), Inches(2.45), Inches(0.34))
        tag.fill.solid()
        tag.fill.fore_color.rgb = COLORS["primary_soft"]
        tag.line.color.rgb = COLORS["primary_soft"]
        tp = tag.text_frame.paragraphs[0]
        tp.text = section
        tp.font.name = FONT
        tp.font.size = Pt(10)
        tp.font.bold = True
        tp.font.color.rgb = COLORS["primary"]
        tp.alignment = PP_ALIGN.CENTER


def add_footer(slide, page, section):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(8.0), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = f"{section} / {page:02d}"
    p.font.name = FONT
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["muted"]


def add_text_box(slide, x, y, w, h, lines, font_size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS["text"]
        p.space_after = Pt(8)
        p.level = 0
    return box


def card(slide, x, y, w, h, title, body=None, accent="primary"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[accent]
    bar.line.color.rgb = COLORS[accent]

    t = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.18), Inches(w - 0.42), Inches(0.3))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]

    if body:
        b = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.58), Inches(w - 0.42), Inches(h - 0.75))
        tf = b.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, line in enumerate(body):
            bp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            bp.text = line
            bp.font.name = FONT
            bp.font.size = Pt(12)
            bp.font.color.rgb = COLORS["muted"]
            bp.space_after = Pt(5)
    return shape


def tag(slide, x, y, text, color="primary"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.55), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["primary_soft"] if color == "primary" else COLORS["white"]
    shape.line.color.rgb = COLORS[color]
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.name = MONO
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLORS[color]
    p.alignment = PP_ALIGN.CENTER
    return shape


def arrow(slide, x1, y1, x2, y2, color="line"):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COLORS[color]
    line.line.width = Pt(2)
    line.line.end_arrowhead = True
    return line


def bullet_slide(prs, page, title, section, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle=subtitle, section=section)
    add_text_box(slide, 0.85, 1.45, 7.4, 4.8, bullets, font_size=19)
    add_footer(slide, page, section)
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def normalize_pptx_zip(path: Path) -> None:
    with zipfile.ZipFile(path, "r") as source:
        entries = [(info, source.read(info.filename)) for info in source.infolist()]

    temp_path = None
    try:
        with tempfile.NamedTemporaryFile(dir=path.parent, suffix=".pptx", delete=False) as temp_file:
            temp_path = Path(temp_file.name)

        with zipfile.ZipFile(temp_path, "w") as target:
            for source_info, data in entries:
                target_info = zipfile.ZipInfo(source_info.filename, ZIP_TIMESTAMP)
                target_info.compress_type = source_info.compress_type
                target_info.comment = source_info.comment
                target_info.extra = source_info.extra
                target_info.external_attr = source_info.external_attr
                target_info.internal_attr = source_info.internal_attr
                target_info.create_system = source_info.create_system
                target_info.create_version = source_info.create_version
                target_info.extract_version = source_info.extract_version
                target_info.flag_bits = source_info.flag_bits
                target_info.volume = source_info.volume
                target.writestr(target_info, data)

        temp_path.replace(path)
    except Exception:
        if temp_path is not None:
            temp_path.unlink(missing_ok=True)
        raise


def main():
    prs = build_presentation()

    # 1. Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "OpenClaw libp2p-mesh 跨实例消息通信机制", "从飞书用户指令到 P2P 网络投递与远端可见消息", "概览")
    card(slide, 0.8, 1.75, 3.0, 1.35, "用户入口", ["飞书中说：给某 instanceId 发消息"], "primary")
    card(slide, 5.0, 1.75, 3.0, 1.35, "P2P 网络", ["libp2p 结构化消息跨实例传输"], "primary")
    card(slide, 9.2, 1.75, 3.0, 1.35, "远端可见", ["投递到远端 Feishu channel"], "success")
    arrow(slide, 3.9, 2.42, 4.85, 2.42)
    arrow(slide, 8.1, 2.42, 9.05, 2.42)
    add_text_box(slide, 1.0, 4.0, 10.8, 1.2, [
        "核心问题：用户指令如何触发 P2P？P2P 消息如何被远端用户看见？",
        "关键机制：instanceId 路由、结构化消息、远端 channel 投递、delivery ACK。"
    ], 19)
    add_footer(slide, 1, "概览")

    # 2. Background
    slide = bullet_slide(prs, 2, "背景与目标：从 peerId 调试到 instanceId 通信", "背景", [
        "peerId 是 libp2p 网络层身份，不适合作为用户操作入口。",
        "instanceId 是 OpenClaw 实例身份，更贴近用户和设备。",
        "目标链路：用户指令 -> Agent 工具 -> P2P 网络 -> 远端用户可见。",
    ])
    card(slide, 8.75, 1.55, 1.9, 1.25, "旧方式", ["peerId 直发", "偏调试"], "warning")
    card(slide, 10.75, 1.55, 1.9, 1.25, "新方式", ["instanceId 发信", "面向用户"], "success")

    # 3. End-to-end flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "总体链路：用户 A 的一句话如何到达用户 B", section="链路")
    steps = [
        ("用户 A 飞书", "指令"),
        ("Agent", "选工具"),
        ("InstanceRouter", "查路由"),
        ("libp2p", "user-message"),
        ("远端 Router", "入站处理"),
        ("Feishu", "sendText"),
        ("用户 B", "看见消息"),
    ]
    x = 0.45
    for i, (name, desc) in enumerate(steps):
        card(slide, x + i * 1.78, 2.05, 1.35, 1.05, name, [desc], "primary" if i < 6 else "success")
        if i < len(steps) - 1:
            arrow(slide, x + i * 1.78 + 1.38, 2.58, x + (i + 1) * 1.78 - 0.08, 2.58)
    tag(slide, 4.95, 3.55, "delivery-ack", "success")
    arrow(slide, 10.9, 3.75, 2.25, 3.75, "success")
    add_text_box(slide, 1.05, 4.55, 11.0, 0.8, ["成功标准不是“P2P 已发出”，而是远端 channel 投递成功并返回 ACK。"], 20)
    add_footer(slide, 3, "链路")

    # 4. Tool triggering
    slide = bullet_slide(prs, 4, "用户指令如何触发 P2P", "工具", [
        "用户给出目标 instanceId 和原始消息内容。",
        "Agent 主路径调用 p2p_send_instance_message({ instanceId, message })。",
        "不再手动先查映射再调用 p2p_send_message。",
        "p2p_send_message 保留为已知 peerId 的低层调试直发。"
    ])
    card(slide, 8.35, 1.45, 4.1, 1.1, "飞书指令", ["给 fhl-enine@... 发送消息：今晚来吃饭"], "primary")
    card(slide, 8.35, 3.0, 4.1, 1.35, "工具调用", ["p2p_send_instance_message", "instanceId + message"], "success")
    arrow(slide, 10.4, 2.6, 10.4, 2.95)

    # 5. Instance route table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "instanceId 如何找到 peerId", section="路由")
    card(slide, 0.8, 1.45, 3.0, 1.25, "mesh 启动 / peer 连接", ["发送 instance-announce"], "primary")
    card(slide, 5.0, 1.45, 3.0, 1.25, "本地路由表", ["instance-peer.json"], "primary")
    card(slide, 9.2, 1.45, 3.0, 1.25, "发送前解析", ["instanceId -> peerId"], "success")
    arrow(slide, 3.85, 2.08, 4.9, 2.08)
    arrow(slide, 8.05, 2.08, 9.1, 2.08)
    add_text_box(slide, 1.0, 3.45, 11.0, 1.6, [
        "映射表由插件自动维护，路径为 ~/.openclaw/libp2p/instance-peer.json。",
        "查询工具：p2p_list_instances / p2p_resolve_instance。",
        "用户不需要在配置中手写 peerId 映射。"
    ], 18)
    add_footer(slide, 5, "路由")

    # 6. Structured messages
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "P2P 网络中传输的是结构化消息", section="消息")
    card(slide, 0.75, 1.55, 3.75, 3.15, "instance-announce", ["路由公告", "instanceId / peerId", "multiaddrs / pubkey"], "primary")
    card(slide, 4.8, 1.55, 3.75, 3.15, "user-message", ["messageId", "fromInstanceId / toInstanceId", "text + reply metadata"], "primary")
    card(slide, 8.85, 1.55, 3.75, 3.15, "delivery-ack", ["ackFor", "ok true/false", "error / inboundChannel"], "success")
    add_text_box(slide, 1.05, 5.35, 11.0, 0.55, ["这三类消息把“发现、发送、确认”拆成清晰协议边界。"], 19)
    add_footer(slide, 6, "消息")

    # 7. Inbound delivery
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "用户如何看见 P2P 消息", section="投递")
    card(slide, 0.8, 1.55, 2.6, 1.25, "远端收到 user-message", ["校验目标 instanceId"], "primary")
    card(slide, 4.0, 1.55, 2.6, 1.25, "读取配置", ["inboundChannel", "inboundTarget"], "primary")
    card(slide, 7.2, 1.55, 2.6, 1.25, "runtime adapter", ["Feishu sendText"], "primary")
    card(slide, 10.4, 1.55, 2.2, 1.25, "用户看见", ["普通文本消息"], "success")
    arrow(slide, 3.45, 2.17, 3.95, 2.17)
    arrow(slide, 6.65, 2.17, 7.15, 2.17)
    arrow(slide, 9.85, 2.17, 10.35, 2.17)
    add_text_box(slide, 1.0, 4.05, 10.8, 1.0, ["当前实现不再 shell out 调 CLI，而是使用 OpenClaw runtime channel outbound adapter 完成投递。"], 19)
    add_footer(slide, 7, "投递")

    # 8. ACK and errors
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "可靠投递：ACK、失败和超时", section="可靠性")
    card(slide, 0.85, 1.45, 3.5, 1.6, "成功", ["远端 channel 投递成功", "返回 ok: true"], "success")
    card(slide, 4.9, 1.45, 3.5, 1.6, "远端失败", ["配置缺失 / 权限失败", "返回 ok: false + error"], "danger")
    card(slide, 8.95, 1.45, 3.5, 1.6, "超时", ["deliveryAckTimeoutMs", "返回 ACK timeout"], "warning")
    add_text_box(slide, 1.0, 4.35, 11.0, 0.9, ["发送方维护 pending ACK map，导师可以把它理解为跨实例投递的闭环确认机制。"], 19)
    add_footer(slide, 8, "可靠性")

    # 9. Tools
    bullet_slide(prs, 9, "新增工具能力", "工具", [
        "低层 peer：p2p_send_message / p2p_broadcast / p2p_list_peers",
        "身份与网络：p2p_get_instance_identity / p2p_get_network_info",
        "实例路由：p2p_list_instances / p2p_resolve_instance / p2p_send_instance_message",
        "普通用户通信优先走 instance 工具，peer 工具主要用于调试。"
    ])

    # 10. Reachability
    slide = bullet_slide(prs, 10, "网络可达性增强", "网络", [
        "mDNS：局域网自动发现。",
        "bootstrap：静态 peer 地址接入。",
        "DHT：WAN peer discovery 与 pubkey registry。",
        "NAT traversal / relay：提升跨 NAT 场景可达性。"
    ])
    card(slide, 8.4, 1.45, 4.0, 2.4, "libp2p 可达性层", ["AutoNAT", "UPnP", "Circuit Relay v2", "DCUtR"], "primary")

    # 11. Safety
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "安全边界与鲁棒性", section="边界")
    card(slide, 0.85, 1.35, 5.45, 3.9, "允许", ["P2P 文本作为普通消息转发", "重复消息复用上次 ACK", "失败原因回传给发送方"], "success")
    card(slide, 7.05, 1.35, 5.45, 3.9, "禁止", ["不把远端文本当系统提示词", "不自动执行远端指令", "不再使用 child_process shell 投递"], "danger")
    add_footer(slide, 11, "边界")

    # 12. Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "总结：三层闭环", section="总结")
    card(slide, 0.85, 1.55, 3.55, 2.35, "用户层", ["只需使用 instanceId", "不理解 peerId 也能发消息"], "primary")
    card(slide, 4.9, 1.55, 3.55, 2.35, "Agent 层", ["主工具清晰", "p2p_send_instance_message"], "primary")
    card(slide, 8.95, 1.55, 3.55, 2.35, "网络层", ["自动路由公告", "结构化消息 + ACK"], "success")
    add_text_box(slide, 1.05, 5.0, 10.8, 0.8, ["后续方向：联系人别名、路由过期清理、受控自动回复、跨 NAT 可视化诊断。"], 18)
    add_footer(slide, 12, "总结")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    normalize_pptx_zip(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
